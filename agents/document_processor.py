"""Document processing utilities for extracting content from various file formats."""

import io
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any
from dataclasses import dataclass

# Document processing libraries
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from PIL import Image
except ImportError:
    Image = None

logger = logging.getLogger(__name__)


@dataclass
class ProcessedDocument:
    """Represents a processed document with extracted content."""
    filename: str
    file_type: str
    content: str
    metadata: Dict[str, Any]
    pages: Optional[int] = None
    word_count: Optional[int] = None
    extracted_images: Optional[List[bytes]] = None
    error: Optional[str] = None


class DocumentProcessor:
    """Handles extraction of content from various document formats."""
    
    SUPPORTED_FORMATS = {
        '.pdf': 'pdf',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint',
        '.docx': 'word',
        '.doc': 'word',
        '.txt': 'text',
        '.md': 'markdown',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.png': 'image',
        '.gif': 'image'
    }
    
    @classmethod
    def process_document(cls, file_path: str, file_content: bytes = None) -> ProcessedDocument:
        """
        Process a document and extract its content.
        
        Args:
            file_path: Path to the document file
            file_content: Optional bytes content if file is already loaded
            
        Returns:
            ProcessedDocument with extracted content
        """
        path = Path(file_path)
        file_ext = path.suffix.lower()
        
        if file_ext not in cls.SUPPORTED_FORMATS:
            return ProcessedDocument(
                filename=path.name,
                file_type='unknown',
                content='',
                metadata={},
                error=f'Unsupported file format: {file_ext}'
            )
        
        file_type = cls.SUPPORTED_FORMATS[file_ext]
        
        # Load file content if not provided
        if file_content is None and path.exists():
            with open(path, 'rb') as f:
                file_content = f.read()
        
        # Process based on file type
        try:
            if file_type == 'pdf':
                return cls._process_pdf(path.name, file_content)
            elif file_type == 'powerpoint':
                return cls._process_powerpoint(path.name, file_content)
            elif file_type == 'word':
                return cls._process_word(path.name, file_content)
            elif file_type in ['text', 'markdown']:
                return cls._process_text(path.name, file_content)
            elif file_type == 'image':
                return cls._process_image(path.name, file_content)
            else:
                return ProcessedDocument(
                    filename=path.name,
                    file_type=file_type,
                    content='',
                    metadata={},
                    error='Processing not implemented for this file type'
                )
        except Exception as e:
            logger.error(f"Error processing {path.name}: {str(e)}")
            return ProcessedDocument(
                filename=path.name,
                file_type=file_type,
                content='',
                metadata={},
                error=str(e)
            )
    
    @classmethod
    def _process_pdf(cls, filename: str, content: bytes) -> ProcessedDocument:
        """Extract text from PDF files."""
        if PyPDF2 is None:
            return ProcessedDocument(
                filename=filename,
                file_type='pdf',
                content='',
                metadata={},
                error='PyPDF2 library not installed'
            )
        
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = []
            metadata = {
                'num_pages': len(pdf_reader.pages),
                'pdf_info': dict(pdf_reader.metadata) if pdf_reader.metadata else {}
            }
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_content.append(f"--- Page {page_num + 1} ---\n{page_text}")
            
            full_text = '\n\n'.join(text_content)
            
            return ProcessedDocument(
                filename=filename,
                file_type='pdf',
                content=full_text,
                metadata=metadata,
                pages=len(pdf_reader.pages),
                word_count=len(full_text.split())
            )
        except Exception as e:
            raise Exception(f"PDF processing error: {str(e)}")
    
    @classmethod
    def _process_powerpoint(cls, filename: str, content: bytes) -> ProcessedDocument:
        """Extract text from PowerPoint presentations."""
        if Presentation is None:
            return ProcessedDocument(
                filename=filename,
                file_type='powerpoint',
                content='',
                metadata={},
                error='python-pptx library not installed'
            )
        
        try:
            ppt_file = io.BytesIO(content)
            presentation = Presentation(ppt_file)
            
            text_content = []
            slide_count = 0
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_count += 1
                slide_text = [f"--- Slide {slide_num + 1} ---"]
                
                # Extract title
                if slide.shapes.title:
                    slide_text.append(f"Title: {slide.shapes.title.text}")
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text:
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))
                
                # Extract notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        slide_text.append(f"Notes: {notes}")
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append('\n'.join(slide_text))
            
            full_text = '\n\n'.join(text_content)
            
            return ProcessedDocument(
                filename=filename,
                file_type='powerpoint',
                content=full_text,
                metadata={'slide_count': slide_count},
                pages=slide_count,
                word_count=len(full_text.split())
            )
        except Exception as e:
            raise Exception(f"PowerPoint processing error: {str(e)}")
    
    @classmethod
    def _process_word(cls, filename: str, content: bytes) -> ProcessedDocument:
        """Extract text from Word documents."""
        if Document is None:
            return ProcessedDocument(
                filename=filename,
                file_type='word',
                content='',
                metadata={},
                error='python-docx library not installed'
            )
        
        try:
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text:
                            row_text.append(cell.text)
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    text_content.append("Table:\n" + '\n'.join(table_text))
            
            full_text = '\n\n'.join(text_content)
            
            # Extract metadata
            core_props = doc.core_properties
            metadata = {
                'author': core_props.author or '',
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'title': core_props.title or '',
                'subject': core_props.subject or ''
            }
            
            return ProcessedDocument(
                filename=filename,
                file_type='word',
                content=full_text,
                metadata=metadata,
                word_count=len(full_text.split())
            )
        except Exception as e:
            raise Exception(f"Word document processing error: {str(e)}")
    
    @classmethod
    def _process_text(cls, filename: str, content: bytes) -> ProcessedDocument:
        """Process plain text or markdown files."""
        try:
            text = content.decode('utf-8', errors='replace')
            
            return ProcessedDocument(
                filename=filename,
                file_type='text',
                content=text,
                metadata={'encoding': 'utf-8'},
                word_count=len(text.split())
            )
        except Exception as e:
            raise Exception(f"Text processing error: {str(e)}")
    
    @classmethod
    def _process_image(cls, filename: str, content: bytes) -> ProcessedDocument:
        """Process image files (extract metadata, prepare for OCR if needed)."""
        if Image is None:
            return ProcessedDocument(
                filename=filename,
                file_type='image',
                content='',
                metadata={},
                error='PIL library not installed'
            )
        
        try:
            img_file = io.BytesIO(content)
            img = Image.open(img_file)
            
            metadata = {
                'format': img.format,
                'mode': img.mode,
                'size': img.size,
                'width': img.width,
                'height': img.height
            }
            
            # Get EXIF data if available
            if hasattr(img, '_getexif') and img._getexif():
                metadata['exif'] = {k: v for k, v in img._getexif().items() if isinstance(v, (str, int, float))}
            
            # Note: Actual OCR would require additional libraries like pytesseract
            content_text = f"Image file: {filename} ({img.width}x{img.height} pixels)"
            
            return ProcessedDocument(
                filename=filename,
                file_type='image',
                content=content_text,
                metadata=metadata,
                extracted_images=[content]
            )
        except Exception as e:
            raise Exception(f"Image processing error: {str(e)}")


class ContentCombiner:
    """Combines content from multiple sources into structured format."""
    
    @staticmethod
    def combine_documents(documents: List[ProcessedDocument]) -> Dict[str, Any]:
        """
        Combine multiple processed documents into a single structure.
        
        Args:
            documents: List of ProcessedDocument objects
            
        Returns:
            Combined content dictionary
        """
        combined = {
            'total_documents': len(documents),
            'total_word_count': 0,
            'document_types': {},
            'all_content': [],
            'errors': [],
            'metadata': {}
        }
        
        for doc in documents:
            # Track document types
            doc_type = doc.file_type
            if doc_type not in combined['document_types']:
                combined['document_types'][doc_type] = 0
            combined['document_types'][doc_type] += 1
            
            # Add content
            if doc.content:
                combined['all_content'].append({
                    'filename': doc.filename,
                    'type': doc.file_type,
                    'content': doc.content,
                    'word_count': doc.word_count
                })
                if doc.word_count:
                    combined['total_word_count'] += doc.word_count
            
            # Track errors
            if doc.error:
                combined['errors'].append({
                    'filename': doc.filename,
                    'error': doc.error
                })
            
            # Merge metadata
            if doc.metadata:
                combined['metadata'][doc.filename] = doc.metadata
        
        return combined
    
    @staticmethod
    def extract_key_sections(content: str) -> Dict[str, str]:
        """
        Extract key sections from combined content.
        
        Args:
            content: Combined text content
            
        Returns:
            Dictionary of extracted sections
        """
        sections = {
            'executive_summary': '',
            'problem_statement': '',
            'solution': '',
            'market_analysis': '',
            'business_model': '',
            'team': '',
            'financials': '',
            'ask': ''
        }
        
        # Simple keyword-based extraction (can be enhanced with ML)
        keywords = {
            'executive_summary': ['executive summary', 'overview', 'about us'],
            'problem_statement': ['problem', 'challenge', 'pain point'],
            'solution': ['solution', 'product', 'service', 'platform'],
            'market_analysis': ['market', 'tam', 'sam', 'industry', 'opportunity'],
            'business_model': ['business model', 'revenue', 'monetization', 'pricing'],
            'team': ['team', 'founders', 'leadership', 'experience'],
            'financials': ['financial', 'revenue', 'projection', 'metrics', 'kpi'],
            'ask': ['ask', 'seeking', 'raise', 'funding', 'investment']
        }
        
        lines = content.lower().split('\n')
        current_section = None
        
        for line in lines:
            # Check if this line indicates a new section
            for section, keys in keywords.items():
                if any(key in line for key in keys):
                    current_section = section
                    break
            
            # Add content to current section
            if current_section and line.strip():
                sections[current_section] += line + '\n'
        
        # Clean up sections
        for key in sections:
            sections[key] = sections[key].strip()
        
        return sections